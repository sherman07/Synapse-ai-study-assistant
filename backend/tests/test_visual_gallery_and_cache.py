import asyncio
import base64
import json
import re
import tempfile
import time
import unittest
from types import SimpleNamespace
from pathlib import Path
from urllib.parse import urlparse
from unittest.mock import patch

from fastapi.testclient import TestClient

from backend.app import app
from backend import app as backend_app_module
from backend.core import config as core_config
from backend.app import analyze_materials
from backend.app import build_analysis_fingerprint
from backend.app import build_visual_gallery
from backend.app import choose_learning_depth
from backend.app import enforce_requested_language
from backend.app import file_to_source_unit
from backend.app import finalize_generated_summary
from backend.app import generate_reference_style_multisource_notes
from backend.app import iter_visual_candidates
from backend.app import link_to_source_unit
from backend.app import rebuild_cached_visual_argument_cards
from backend.app import render_pptx_source_preview_svg_images
from backend.app import source_preview_image_url
from backend.app import validate_source_strict_summary
from backend.core.analysis_cache import AnalysisCacheStore
from backend.core import visual_assets
from backend.core.visual_assets import prune_unavailable_visual_markers, visual_asset_url_for_browser

try:
    from pptx import Presentation
except Exception:
    Presentation = None


def assert_served_visual_asset(test_case, url: str, expected_content_type: str = "image/png"):
    parsed = urlparse(url)
    test_case.assertEqual(parsed.scheme, "http")
    test_case.assertEqual(parsed.netloc, "127.0.0.1:8001")
    test_case.assertTrue(parsed.path.startswith("/assets/visuals/"))

    response = TestClient(app).get(parsed.path)
    test_case.assertEqual(response.status_code, 200)
    test_case.assertEqual(response.headers.get("content-type", "").split(";")[0], expected_content_type)


class ApiShapeTests(unittest.TestCase):
    def test_liveness_endpoint_is_minimal_for_platform_health_checks(self):
        response = TestClient(app).get("/healthz")

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.text, "ok")
        self.assertLess(len(response.content), 32)

    def test_health_response_shape(self):
        response = TestClient(app).get("/health")

        self.assertEqual(response.status_code, 200)
        payload = response.json()
        self.assertEqual(payload["status"], "ok")
        self.assertIn("analysis_model", payload)
        self.assertIn("fallback_model", payload)
        self.assertIn("mindmap_model", payload)
        self.assertIn("title_model", payload)
        self.assertIn("visual_image_guide_model", payload)
        self.assertIn("openai_timeout_seconds", payload)
        self.assertIn("pdf_visual_extraction_enabled", payload)
        self.assertIn("embedded_youtube_sources_enabled", payload)
        self.assertIn("youtube_ytdlp_fallback_enabled", payload)
        self.assertIn("public_backend_base_url", payload)
        self.assertIn("supabase_auth_configured", payload)
        self.assertIn("synapse_email_delivery_configured", payload)
        self.assertIn("runtime_assets_dir", payload)
        self.assertNotIn("OPENAI_API_KEY", payload)

    def test_default_frontend_base_url_matches_documented_vite_server(self):
        self.assertEqual(
            backend_app_module.SYNAPSE_FRONTEND_BASE_URL,
            "http://127.0.0.1:5175/frontend",
        )

    def test_default_cors_allows_documented_static_frontend_server(self):
        self.assertTrue(
            {
                "http://127.0.0.1:5500",
                "http://localhost:5500",
                "http://127.0.0.1:5176",
                "http://localhost:5176",
            }.issubset(set(backend_app_module.CORS_ALLOW_ORIGINS))
        )

    def test_default_cors_allows_private_lan_vite_frontend(self):
        response = TestClient(app).options(
            "/health",
            headers={
                "Origin": "http://192.168.1.141:5176",
                "Access-Control-Request-Method": "GET",
            },
        )

        self.assertEqual(response.status_code, 200)
        self.assertEqual(
            response.headers.get("access-control-allow-origin"),
            "http://192.168.1.141:5176",
        )

    def test_backend_loads_gemini_settings_from_separate_env_file(self):
        self.assertIn(
            backend_app_module.BACKEND_PACKAGE_DIR / ".env.gemini",
            core_config.GEMINI_ENV_PATHS,
        )
        self.assertIn(
            backend_app_module.BACKEND_PACKAGE_DIR / ".env.gemini",
            core_config.CONFIG_ENV_PATHS,
        )

    def test_backend_loads_gpt_settings_from_separate_env_file(self):
        self.assertIn(
            backend_app_module.BACKEND_PACKAGE_DIR / ".env.gpt",
            core_config.GPT_ENV_PATHS,
        )
        self.assertIn(
            backend_app_module.BACKEND_PACKAGE_DIR / "core" / ".env.gpt",
            core_config.GPT_ENV_PATHS,
        )
        self.assertIn(
            backend_app_module.BACKEND_PACKAGE_DIR / "core" / ".env.gpt",
            core_config.CONFIG_ENV_PATHS,
        )

    def test_gpt_env_values_replace_placeholder_openai_settings(self):
        environ = {
            "OPENAI_API_KEY": "__ADD_YOUR_OPENAI_API_KEY__",
            "OPENAI_REALTIME_MODEL": "__ADD_YOUR_REALTIME_MODEL__",
            "UNRELATED_VALUE": "keep-me",
        }

        core_config.apply_env_values(
            {
                "OPENAI_API_KEY": "sk-proj-real-test-key",
                "OPENAI_REALTIME_MODEL": "gpt-realtime-2",
                "UNRELATED_VALUE": "replace-me",
            },
            environ=environ,
            override_placeholders=True,
        )

        self.assertEqual(environ["OPENAI_API_KEY"], "sk-proj-real-test-key")
        self.assertEqual(environ["OPENAI_REALTIME_MODEL"], "gpt-realtime-2")
        self.assertEqual(environ["UNRELATED_VALUE"], "keep-me")

    def test_model_for_depth_uses_gemini_models_when_text_provider_is_gemini(self):
        with (
            patch.object(core_config, "AI_TEXT_PROVIDER", "gemini"),
            patch.object(core_config, "GEMINI_AUTH_MODE", "api_key", create=True),
            patch.object(core_config, "GEMINI_FOCUSED_MODEL", "gemini-focused"),
            patch.object(core_config, "GEMINI_STANDARD_MODEL", "gemini-standard"),
            patch.object(core_config, "GEMINI_DETAILED_MODEL", "gemini-detailed"),
            patch.object(core_config, "GEMINI_COMPREHENSIVE_MODEL", "gemini-comprehensive"),
        ):
            self.assertEqual(core_config.model_for_depth("focused"), "gemini-focused")
            self.assertEqual(core_config.model_for_depth("standard"), "gemini-standard")
            self.assertEqual(core_config.model_for_depth("detailed"), "gemini-detailed")
            self.assertEqual(core_config.model_for_depth("comprehensive"), "gemini-comprehensive")

    def test_unconfigured_gemini_request_falls_back_to_openai(self):
        with (
            patch.object(core_config, "AI_TEXT_PROVIDER", "openai"),
            patch.object(core_config, "OPENAI_API_KEY", "test-openai-key"),
            patch.object(core_config, "client", object()),
            patch.object(core_config, "GEMINI_AUTH_MODE", "adc"),
            patch.object(core_config, "GEMINI_PROJECT_ID", ""),
        ):
            token = core_config.set_request_text_provider("gemini")
            try:
                self.assertEqual(
                    core_config.active_text_provider(),
                    "openai",
                    "an unavailable optional Gemini selection must not break analysis",
                )
            finally:
                core_config.reset_request_text_provider(token)

    def test_generate_chat_uses_gemini_client_when_provider_is_gemini(self):
        class FakeCompletions:
            def __init__(self, content):
                self.content = content
                self.calls = []

            def create(self, **kwargs):
                self.calls.append(kwargs)
                return SimpleNamespace(
                    choices=[
                        SimpleNamespace(
                            message=SimpleNamespace(content=self.content),
                        ),
                    ],
                    usage=SimpleNamespace(prompt_tokens=1, completion_tokens=2, total_tokens=3),
                )

        class FakeClient:
            def __init__(self, content):
                self.completions = FakeCompletions(content)
                self.chat = SimpleNamespace(completions=self.completions)

        openai_client = FakeClient("openai output")
        gemini_client = FakeClient("gemini output")

        trace_token = backend_app_module.begin_ai_call_trace()
        with (
            patch.object(backend_app_module, "AI_TEXT_PROVIDER", "gemini", create=True),
            patch.object(backend_app_module, "client", openai_client),
            patch.object(backend_app_module, "text_generation_client", return_value=gemini_client, create=True),
        ):
            result = backend_app_module.generate_chat(
                [{"role": "user", "content": "Use the selected Synapse prompt."}],
                model="gemini-2.5-flash",
                temperature=0,
                max_tokens=20,
            )
        trace = backend_app_module.current_ai_call_trace()
        backend_app_module.reset_ai_call_trace(trace_token)

        self.assertEqual(result, "gemini output")
        self.assertEqual(len(gemini_client.completions.calls), 1)
        self.assertEqual(len(openai_client.completions.calls), 0)
        self.assertEqual(gemini_client.completions.calls[0]["model"], "gemini-2.5-flash")
        self.assertEqual(len(trace), 1)
        self.assertEqual(trace[0]["status"], "success")
        self.assertTrue(trace[0]["api_request_attempted"])
        self.assertEqual(trace[0]["model"], "gemini-2.5-flash")
        self.assertEqual(trace[0]["total_tokens"], 3)

    def test_gemini_adc_client_uses_vertex_openai_endpoint_and_token(self):
        class FakeCredentials:
            token = "adc-token"

            def __init__(self):
                self.refreshed = False

            def refresh(self, request):
                self.refreshed = request == "adc-request"

        credentials = FakeCredentials()

        with (
            patch.object(core_config, "AI_TEXT_PROVIDER", "gemini"),
            patch.object(core_config, "GEMINI_AUTH_MODE", "adc", create=True),
            patch.object(core_config, "GEMINI_PROJECT_ID", "synapse-project", create=True),
            patch.object(core_config, "GEMINI_LOCATION", "us-central1", create=True),
            patch.object(core_config, "google_auth_default", return_value=(credentials, None), create=True),
            patch.object(core_config, "google_auth_request", return_value="adc-request", create=True),
            patch.object(core_config, "OpenAI") as openai_mock,
        ):
            client = core_config.text_generation_client()

        self.assertIs(client, openai_mock.return_value)
        self.assertTrue(credentials.refreshed)
        openai_mock.assert_called_once_with(
            api_key="adc-token",
            base_url="https://us-central1-aiplatform.googleapis.com/v1/projects/synapse-project/locations/us-central1/endpoints/openapi",
            timeout=core_config.OPENAI_TIMEOUT_SECONDS,
        )

    def test_gemini_adc_global_location_uses_global_vertex_endpoint(self):
        with (
            patch.object(core_config, "GEMINI_PROJECT_ID", "synapse-project", create=True),
            patch.object(core_config, "GEMINI_LOCATION", "global", create=True),
        ):
            self.assertEqual(
                core_config.gemini_vertex_openai_base_url(),
                "https://aiplatform.googleapis.com/v1beta1/projects/synapse-project/locations/global/endpoints/openapi",
            )

    def test_explicit_detail_level_overrides_auto_depth(self):
        payload = choose_learning_depth("Tiny note about slope.", [], "detailed")

        self.assertEqual(payload["depth"], "detailed")
        self.assertTrue(payload["override"])
        self.assertEqual(payload["requested_detail_level"], "detailed")
        self.assertEqual(payload["auto_selected_depth"], "focused")

    def test_auto_detail_level_keeps_adaptive_depth(self):
        payload = choose_learning_depth("Tiny note about slope.", [], "auto")

        self.assertEqual(payload["depth"], "focused")
        self.assertFalse(payload["override"])
        self.assertEqual(payload["requested_detail_level"], "auto")

    def test_analysis_deadline_skips_optional_model_stages(self):
        with (
            patch("backend.app.ANALYSIS_MAX_SECONDS", 60),
            patch("backend.app.analysis_elapsed_seconds_since", return_value=60.0),
            patch("backend.app.should_run_optional_analysis_stage", return_value=False),
            patch("backend.app.require_text_ai"),
            patch("backend.app.cache_get", return_value=None),
            patch("backend.app.cache_set"),
            patch(
                "backend.app.generate_reference_style_multisource_notes",
                return_value="# Overview\n\nTiny source note about slope.\n\n## Key Ideas\n\nSlope compares rise and run.",
            ),
            patch("backend.app.make_notes_title", side_effect=AssertionError("title model should be skipped")),
            patch("backend.app.generate_ai_mind_map", side_effect=AssertionError("AI mind map should be skipped")),
        ):
            payload = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="Tiny note about slope.",
                preferred_language="english",
                detail_level="auto",
                prompt_mode="professor_mode",
                client_fingerprint="",
            ))

        self.assertNotIn("error", payload)
        self.assertIn("title", payload["optional_stages_skipped"])
        self.assertIn("mind_map", payload["optional_stages_skipped"])
        self.assertEqual(payload["analysis_max_seconds"], 60)
        self.assertGreaterEqual(payload["analysis_elapsed_seconds"], 60)
        self.assertEqual(payload["language"], "english")
        self.assertEqual(payload["output_language"], "english")
        self.assertTrue(payload["mind_map"].get("branches"))

    def test_analyze_bounds_optional_post_note_model_timeouts(self):
        captured_timeouts = {}

        def fake_language_rewrite(summary, preferred_language, **kwargs):
            captured_timeouts["language"] = kwargs.get("request_timeout")
            return summary

        def fake_title(summary, candidates, **kwargs):
            captured_timeouts["title"] = kwargs.get("request_timeout")
            return "Bounded Optional Title"

        def fake_localise_title(title, preferred_language, **kwargs):
            captured_timeouts["localise_title"] = kwargs.get("request_timeout")
            return title

        def fake_mind_map(title, section_map, preferred_language, depth, prompt_mode, **kwargs):
            captured_timeouts["mind_map"] = kwargs.get("request_timeout")
            return {"center": title, "branches": [{"label": "Key Ideas", "points": []}]}

        with (
            patch("backend.app.require_text_ai"),
            patch("backend.app.cache_get", return_value=None),
            patch("backend.app.cache_set"),
            patch("backend.app.should_run_optional_analysis_stage", return_value=True),
            patch("backend.app.analysis_model_call_timeout", side_effect=[33.0, 12.0, 11.0, 22.0]),
            patch(
                "backend.app.generate_reference_style_multisource_notes",
                return_value="# Overview\n\nTiny source note about slope.\n\n## Key Ideas\n\nSlope compares rise and run.",
            ),
            patch("backend.app.enforce_requested_language", side_effect=fake_language_rewrite),
            patch("backend.app.make_notes_title", side_effect=fake_title),
            patch("backend.app.localise_title_if_needed", side_effect=fake_localise_title),
            patch("backend.app.generate_ai_mind_map", side_effect=fake_mind_map),
        ):
            payload = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="Tiny note about slope.",
                preferred_language="english",
                detail_level="auto",
                prompt_mode="professor_mode",
                client_fingerprint="",
            ))

        self.assertNotIn("error", payload)
        self.assertEqual(captured_timeouts["language"], 33.0)
        self.assertEqual(captured_timeouts["title"], 12.0)
        self.assertEqual(captured_timeouts["localise_title"], 11.0)
        self.assertEqual(captured_timeouts["mind_map"], 22.0)

    def test_analyze_route_returns_http_error_for_empty_request(self):
        with patch("backend.app.require_text_ai"):
            response = TestClient(app).post("/analyze", data={})

        self.assertEqual(response.status_code, 400)
        self.assertEqual(response.json(), {"error": "No readable files, links, or text were provided."})

    def test_analyze_route_returns_service_error_when_text_provider_missing(self):
        with patch("backend.app.require_text_ai", side_effect=RuntimeError("GEMINI_API_KEY is missing.")):
            response = TestClient(app).post("/analyze", data={"free_text": "short source"})

        self.assertEqual(response.status_code, 503)
        self.assertEqual(response.json(), {"error": "GEMINI_API_KEY is missing."})

    def test_analyze_response_exposes_source_identity_alias(self):
        with (
            patch("backend.app.require_text_ai"),
            patch("backend.app.cache_get", return_value=None),
            patch("backend.app.cache_set"),
            patch("backend.app.should_run_optional_analysis_stage", return_value=False),
            patch("backend.app.build_visual_gallery", return_value=[]),
            patch(
                "backend.app.generate_reference_style_multisource_notes",
                return_value="# Overview\n\nThis source explains evidence and revision use.",
            ),
        ):
            payload = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="A short source about evidence, concepts, and revision use.",
                preferred_language="english",
                detail_level="auto",
                prompt_mode="professor_mode",
                client_fingerprint="",
            ))

        self.assertNotIn("error", payload)
        self.assertTrue(payload["source_identity"].startswith("text:"))
        self.assertEqual(payload["source_identity"], payload["primary_source_identity"])

    def test_cached_result_visual_rebuild_never_calls_model(self):
        visual_parts = [
            {
                "type": "text",
                "text": (
                    "IN-TEXT SOURCE FIGURE FROM cached lecture — PDF page 2. "
                    "Page text preview: table graph data results comparison. "
                    "Image-count=1; drawing-count=12; visual-score=28."
                ),
            },
            {"type": "image_url", "image_url": {"url": "data:image/png;base64,AA=="}},
        ]
        source_meta = {
            "display_name": "cached lecture",
            "source_identity": "url:https://example.com/cached-lecture",
            "title_candidate": "Cached lecture",
            "content_hash": "cached-lecture",
            "text_excerpt": "The source includes a result table used as evidence.",
            "visual_parts": visual_parts,
            "url": "https://example.com/cached-lecture",
        }
        cached_result = {
            "title": "Cached lecture notes",
            "summary": "# Overview\n\nCached notes discuss the source result table.",
            "sections": {"Overview": "Cached notes discuss the source result table."},
            "connections": [],
            "mind_map": {"center": "Cached lecture notes", "branches": []},
            "visual_gallery": [],
            "primary_source_identity": source_meta["source_identity"],
        }

        with (
            patch("backend.app.require_text_ai"),
            patch("backend.app.cache_get", return_value=cached_result),
            patch("backend.app.rebuild_cached_visual_argument_cards", return_value=[]),
            patch("backend.app.link_to_source_unit", return_value=(
                [{"type": "text", "text": "Cached lecture with a result table."}],
                source_meta,
            )),
            patch("backend.app.generate_chat", side_effect=AssertionError("cache hit should not call the model")),
        ):
            payload = asyncio.run(analyze_materials(
                files=[],
                links=json.dumps(["https://example.com/cached-lecture"]),
                free_text="",
                preferred_language="english",
                detail_level="auto",
                prompt_mode="professor_mode",
                client_fingerprint="",
            ))

        self.assertNotIn("error", payload)
        self.assertTrue(payload["cached"])
        self.assertEqual(payload["ai_generation_source"], "cache")
        self.assertEqual(payload["ai_model_call_count"], 0)
        self.assertFalse(payload["ai_fallback_used"])
        self.assertNotIn("[[VISUAL:", payload["summary"])
        self.assertNotIn("This source figure belongs in the notes", payload["summary"])
        self.assertNotIn("Figure focus:", payload["summary"])
        self.assertEqual(len(payload["visual_gallery"]), 1)
        self.assertEqual(payload["source_evidence_cards"], payload["visual_gallery"])
        self.assertEqual(payload["figure_cards"], payload["visual_gallery"])

    def test_analyze_reports_fallback_when_main_note_model_fails(self):
        with (
            patch("backend.app.require_text_ai"),
            patch("backend.app.cache_get", return_value=None),
            patch("backend.app.cache_set"),
            patch("backend.app.should_run_optional_analysis_stage", return_value=False),
            patch("backend.app.build_visual_gallery", return_value=[]),
            patch("backend.app.persist_generated_analysis_result", return_value=None),
            patch("backend.app.generate_chat", side_effect=RuntimeError("forced model failure")),
        ):
            payload = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="Human nature lecture about aggression, cooperation, evidence, and exam interpretation.",
                preferred_language="english",
                detail_level="auto",
                prompt_mode="professor_mode",
                client_fingerprint="",
            ))

        self.assertNotIn("error", payload)
        self.assertFalse(payload["cached"])
        self.assertEqual(payload["ai_generation_source"], "fallback")
        self.assertTrue(payload["ai_fallback_used"])
        self.assertIn("main_notes", payload["ai_generation"]["fallback_stages"])
        self.assertIn("forced model failure", payload["ai_generation"]["last_error"])
        self.assertIn("Professional Study Guide", payload["summary"])

    def test_deadline_skips_visual_filter_and_note_expansion_model_stages(self):
        long_summary = "# Overview\n\n" + ("This source explains table data, comparison, evidence, and exam use. " * 700)
        candidate = {
            "source_index": 1,
            "source_title": "Lecture results",
            "display_name": "lecture.pdf",
            "caption": "A table with data, results, comparison groups, and statistical evidence.",
            "location": "PDF page 2",
            "url": "data:image/png;base64,AA==",
            "score": 32,
            "visual_kind": "data/table",
            "teaching_signals": 4,
            "decorative_signals": 0,
            "is_likely_decorative": False,
        }
        source_units = [{
            "display_name": "lecture.pdf",
            "title_candidate": "Lecture results",
            "text_excerpt": "The source includes a table with data, results, and comparison evidence.",
            "visual_parts": [],
        }]
        skipped = []

        with (
            patch("backend.app.analysis_remaining_seconds_since", return_value=20),
            patch("backend.app.select_visual_candidates_for_argument", return_value=[candidate]),
            patch("backend.app.generate_chat", return_value=long_summary) as generate_chat_mock,
            patch("backend.app.expand_sparse_inline_summary", side_effect=AssertionError("expansion should be skipped")),
        ):
            summary = generate_reference_style_multisource_notes(
                source_units,
                "english",
                {"depth": "detailed", "config": backend_app_module.DEPTH_CONFIG["detailed"]},
                "professor_mode",
                analysis_started_at=1.0,
                skipped_optional_stages=skipped,
            )

        self.assertEqual(generate_chat_mock.call_count, 1)
        self.assertIn("visual_card_filter", skipped)
        self.assertIn("note_expansion", skipped)
        self.assertIn("[[VISUAL:0]]", summary)
        self.assertNotIn("This source figure belongs in the notes", summary)
        self.assertNotIn("Figure focus:", summary)

    def test_analyze_cache_preserves_browser_visual_metadata(self):
        captured = {}
        gallery = [{
            "index": 0,
            "url": "http://127.0.0.1:8001/assets/visuals/result-table.png",
            "title": "Result table",
            "caption": "Table with data and results.",
            "visual_kind": "data/table",
        }]

        def capture_cache_set(fingerprint, result):
            captured["fingerprint"] = fingerprint
            captured["result"] = result

        with (
            patch("backend.app.require_text_ai"),
            patch("backend.app.cache_get", return_value=None),
            patch("backend.app.cache_set", side_effect=capture_cache_set),
            patch("backend.app.should_run_optional_analysis_stage", return_value=False),
            patch("backend.app.build_visual_gallery", return_value=gallery),
            patch(
                "backend.app.generate_reference_style_multisource_notes",
                return_value="# Overview\n\nThe table evidence is explained near [[VISUAL:0]].\n\n## Key Ideas\n\nData supports the claim.",
            ),
        ):
            payload = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="Tiny note about a result table.",
                preferred_language="english",
                detail_level="auto",
                prompt_mode="professor_mode",
                client_fingerprint="",
            ))

        self.assertNotIn("error", payload)
        self.assertIn("[[VISUAL:0]]", payload["summary"])
        self.assertEqual(payload["visual_gallery"], gallery)
        self.assertEqual(payload["visuals"], gallery)
        self.assertEqual(payload["source_evidence_cards"], gallery)
        self.assertEqual(payload["figure_cards"], gallery)
        self.assertIn("[[VISUAL:0]]", captured["result"]["summary"])
        self.assertEqual(captured["result"]["visual_gallery"], gallery)
        self.assertEqual(captured["result"]["visuals"], gallery)
        self.assertEqual(captured["result"]["source_evidence_cards"], gallery)
        self.assertEqual(captured["result"]["figure_cards"], gallery)

    def test_cached_result_uses_cached_visual_gallery_when_live_rebuild_is_empty(self):
        asset_path = backend_app_module.RUNTIME_ASSETS_DIR / "visuals" / "test-cached-table.png"
        asset_path.parent.mkdir(parents=True, exist_ok=True)
        asset_path.write_bytes(b"cached visual")
        cached_gallery = [{
            "index": 0,
            "url": "http://127.0.0.1:8001/assets/visuals/test-cached-table.png",
            "title": "Cached result table",
            "caption": "Cached table metadata.",
            "visual_kind": "data/table",
        }]
        cached_result = {
            "title": "Cached notes",
            "summary": "# Overview\n\nCached notes keep [[VISUAL:0]].",
            "sections": {"Overview": "Cached notes keep [[VISUAL:0]]."},
            "connections": [],
            "mind_map": {"center": "Cached notes", "branches": []},
            "visual_gallery": cached_gallery,
            "primary_source_identity": "text:cached",
        }

        try:
            with (
                patch("backend.app.require_text_ai"),
                patch("backend.app.cache_get", return_value=cached_result),
                patch("backend.app.rebuild_cached_visual_argument_cards", return_value=[]),
                patch("backend.app.build_visual_gallery", return_value=[]),
                patch("backend.app.generate_chat", side_effect=AssertionError("cache hit should not need a model call")),
            ):
                payload = asyncio.run(analyze_materials(
                    files=[],
                    links="[]",
                    free_text="Tiny note about a result table.",
                    preferred_language="english",
                    detail_level="auto",
                    prompt_mode="professor_mode",
                    client_fingerprint="",
                ))
        finally:
            asset_path.unlink(missing_ok=True)

        self.assertTrue(payload["cached"])
        self.assertIn("[[VISUAL:0]]", payload["summary"])
        self.assertEqual(payload["visual_gallery"], cached_gallery)
        self.assertEqual(payload["visuals"], cached_gallery)
        self.assertEqual(payload["source_evidence_cards"], cached_gallery)
        self.assertEqual(payload["figure_cards"], cached_gallery)

    def test_cached_result_does_not_return_missing_runtime_visual_asset(self):
        cached_result = {
            "title": "Cached notes",
            "summary": "# Overview\n\nCached notes still reference [[VISUAL:0]].",
            "sections": {"Overview": "Cached notes still reference [[VISUAL:0]]."},
            "connections": [],
            "mind_map": {"center": "Cached notes", "branches": []},
            "visual_gallery": [{
                "index": 0,
                "url": "http://127.0.0.1:8001/assets/visuals/definitely-missing-cached-asset.png",
                "title": "Missing cached table",
                "caption": "The cache entry points to a runtime asset that no longer exists.",
                "visual_kind": "data/table",
            }],
            "primary_source_identity": "text:cached",
        }

        with (
            patch("backend.app.require_text_ai"),
            patch("backend.app.cache_get", return_value=cached_result),
            patch("backend.app.rebuild_cached_visual_argument_cards", return_value=[]),
            patch("backend.app.build_visual_gallery", return_value=[]),
            patch("backend.app.generate_chat", side_effect=AssertionError("cache hit should not need a model call")),
        ):
            payload = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="Tiny note about a missing cached result table.",
                preferred_language="english",
                detail_level="auto",
                prompt_mode="professor_mode",
                client_fingerprint="",
            ))

        self.assertTrue(payload["cached"])
        self.assertEqual(payload["visual_gallery"], [])
        self.assertEqual(payload["visuals"], [])
        self.assertNotIn("[[VISUAL:0]]", payload["summary"])

    def test_language_enforcement_skips_rewrite_when_english_already_satisfied(self):
        summary = "# Overview\n\nThis study note explains the source evidence, comparison table, and revision use."

        with patch("backend.app.generate_chat", side_effect=AssertionError("English notes should not be rewritten")):
            result = enforce_requested_language(summary, "english")

        self.assertEqual(result, summary)

    def test_language_enforcement_skips_rewrite_when_chinese_already_satisfied(self):
        summary = "# 概述\n\n这份学习笔记解释来源证据、比较表格、核心概念和复习用途，帮助学生直接按照材料复习。"

        with patch("backend.app.generate_chat", side_effect=AssertionError("Chinese notes should not be rewritten")):
            result = enforce_requested_language(summary, "simplified_chinese")

        self.assertEqual(result, summary)

    def test_language_enforcement_rewrites_visible_mismatch(self):
        rewritten = "# 概述\n\n这份笔记已经改写为中文。"

        with patch("backend.app.generate_chat", return_value=rewritten) as generate_chat_mock:
            result = enforce_requested_language("# Overview\n\nThis note is still English.", "simplified_chinese")

        self.assertEqual(result, rewritten)
        self.assertEqual(generate_chat_mock.call_count, 1)

    def test_language_enforcement_rewrites_simplified_when_traditional_requested(self):
        rewritten = "# 概覽\n\n這份筆記已經改寫為繁體中文。"

        with patch("backend.app.generate_chat", return_value=rewritten) as generate_chat_mock:
            result = enforce_requested_language("# 概述\n\n这份学习笔记解释来源证据和复习用途。", "traditional_chinese")

        self.assertEqual(result, rewritten)
        self.assertEqual(generate_chat_mock.call_count, 1)


class EmbeddedYoutubeSafetyTests(unittest.TestCase):
    def test_uploaded_file_does_not_expand_embedded_youtube_when_disabled(self):
        previous_value = getattr(backend_app_module, "ENABLE_EMBEDDED_YOUTUBE_SOURCES", None)
        backend_app_module.ENABLE_EMBEDDED_YOUTUBE_SOURCES = False
        try:
            with patch("backend.app.link_to_source_unit") as link_to_source_unit:
                result = backend_app_module.expand_embedded_youtube_sources(
                    "Watch https://www.youtube.com/watch?v=dQw4w9WgXcQ for more context.",
                    {"display_name": "Lecture handout.pdf"},
                    set(),
                )

            self.assertEqual(result, ([], [], []))
            link_to_source_unit.assert_not_called()
        finally:
            if previous_value is None:
                delattr(backend_app_module, "ENABLE_EMBEDDED_YOUTUBE_SOURCES")
            else:
                backend_app_module.ENABLE_EMBEDDED_YOUTUBE_SOURCES = previous_value

    def test_ytdlp_caption_fallback_is_skipped_when_disabled(self):
        previous_value = getattr(backend_app_module, "ENABLE_YOUTUBE_YTDLP_FALLBACK", None)
        backend_app_module.ENABLE_YOUTUBE_YTDLP_FALLBACK = False
        try:
            with patch.object(backend_app_module, "yt_dlp") as yt_dlp:
                transcript = backend_app_module.fetch_youtube_subtitle_transcript(
                    "https://www.youtube.com/watch?v=dQw4w9WgXcQ"
                )

            self.assertEqual(transcript, "")
            yt_dlp.YoutubeDL.assert_not_called()
        finally:
            if previous_value is None:
                delattr(backend_app_module, "ENABLE_YOUTUBE_YTDLP_FALLBACK")
            else:
                backend_app_module.ENABLE_YOUTUBE_YTDLP_FALLBACK = previous_value


class SourceStrictNotesTests(unittest.TestCase):
    def test_source_strict_validation_rebuilds_required_sections_and_dedupes_repeat(self):
        summary = """
# Vaccination ethics

## Direct Source Claims

- Mandates must be necessary and proportionate (Slide 15).
- Mandates must be necessary and proportionate (Slide 15).

## Source Evidence

The lecture uses Jacobson v. Massachusetts to illustrate necessity and proportionality (Slide 21).
The lecture uses Jacobson v. Massachusetts to illustrate necessity and proportionality (Slide 21).
""".strip()

        validated = validate_source_strict_summary(summary, "english", "standard_notes")

        self.assertIn("## Source Question", validated)
        self.assertIn("## Direct Source Claims", validated)
        self.assertIn("## Source Evidence", validated)
        self.assertIn("## Inferences Allowed By The Source", validated)
        self.assertIn("## Gaps / Limits", validated)
        self.assertIn("## Exam / Research Use", validated)
        self.assertIn("## Compact Revision Summary", validated)
        self.assertIn("Not enough evidence from the uploaded source.", validated)
        self.assertNotIn("[Direct from source]", validated)
        self.assertNotIn("[Tutor explanation]", validated)
        self.assertEqual(validated.count("Mandates must be necessary and proportionate (Slide 15)."), 1)
        self.assertEqual(
            validated.count("Jacobson v. Massachusetts to illustrate necessity and proportionality (Slide 21)."),
            1,
        )

    def test_source_strict_finalize_respects_quick_review_word_limit(self):
        repeated_sentence = "Necessity, proportionality, and liberty must be balanced with cited lecture evidence (Slide 15)."
        summary = "# Vaccination policy\n\n## Direct Source Claims\n\n" + "\n\n".join(repeated_sentence for _ in range(80))

        final = finalize_generated_summary(
            summary,
            requested_language="english",
            generation_language="english",
            prompt_mode="source_strict_research_mode",
            note_length_mode="quick_review",
            attach_visuals=False,
        )

        word_count = len(re.findall(r"[\u4e00-\u9fff]|\b[\w'-]+\b", final))
        self.assertLessEqual(word_count, 560)
        self.assertIn("## Direct Source Claims", final)

    def test_analysis_fingerprint_changes_with_note_length_mode(self):
        units = [{
            "source_identity": "text:abc",
            "content_hash": "hash-1",
        }]

        quick = build_analysis_fingerprint("english", units, "auto", "source_strict_research_mode", "quick_review")
        deep = build_analysis_fingerprint("english", units, "auto", "source_strict_research_mode", "deep_study")

        self.assertNotEqual(quick, deep)

    def test_analysis_fingerprint_changes_with_ai_provider(self):
        units = [{
            "source_identity": "text:abc",
            "content_hash": "hash-1",
        }]

        openai = build_analysis_fingerprint(
            "english",
            units,
            "auto",
            "professor_mode",
            "standard_notes",
            ai_provider="openai",
        )
        gemini = build_analysis_fingerprint(
            "english",
            units,
            "auto",
            "professor_mode",
            "standard_notes",
            ai_provider="gemini",
        )

        self.assertNotEqual(openai, gemini)

    def test_analysis_fingerprint_changes_with_visual_pipeline_version(self):
        units = [{
            "source_identity": "text:abc",
            "content_hash": "hash-1",
        }]

        with patch.object(backend_app_module, "VISUAL_PIPELINE_VERSION", "no-inline-visuals-v2", create=True):
            first = build_analysis_fingerprint("english", units, "auto", "professor_mode", "standard_notes")
        with patch.object(backend_app_module, "VISUAL_PIPELINE_VERSION", "no-inline-visuals-v3", create=True):
            second = build_analysis_fingerprint("english", units, "auto", "professor_mode", "standard_notes")

        self.assertNotEqual(first, second)

    def test_analyze_stores_raw_and_display_summaries_separately(self):
        captured = {}
        gallery = [{
            "index": 0,
            "url": "http://127.0.0.1:8001/assets/visuals/result-table.png",
            "title": "Result table",
            "caption": "Table with data and results.",
            "visual_kind": "data/table",
        }]

        def capture_cache_set(fingerprint, result):
            captured["result"] = result

        with (
            patch("backend.app.require_text_ai"),
            patch("backend.app.cache_get", return_value=None),
            patch("backend.app.cache_set", side_effect=capture_cache_set),
            patch("backend.app.should_run_optional_analysis_stage", return_value=False),
            patch("backend.app.build_visual_gallery", return_value=gallery),
            patch(
                "backend.app.generate_reference_style_multisource_notes",
                return_value=(
                    "# Overview\n\n"
                    "The lecture explains the result table as source evidence near [[VISUAL:0]].\n\n"
                    "## Source Examples and Evidence\n\n"
                    "This source figure belongs in the notes because it shows the table.\n\n"
                    "## Core Argument\n\n"
                    "The result table matters because it limits the claim."
                ),
            ),
        ):
            payload = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="Tiny note about a result table.",
                preferred_language="english",
                detail_level="auto",
                prompt_mode="professor_mode",
                client_fingerprint="",
            ))

        self.assertNotIn("error", payload)
        self.assertEqual(payload["summary"], payload["display_summary"])
        self.assertEqual(captured["result"]["summary"], captured["result"]["display_summary"])
        self.assertIn("raw_summary", payload)
        self.assertIn("display_summary", payload)
        for key in ("summary", "display_summary", "raw_summary"):
            self.assertIn("[[VISUAL:0]]", payload[key])
            self.assertNotIn("Source Examples and Evidence", payload[key])
            self.assertNotIn("This source figure belongs in the notes", payload[key])
        self.assertEqual(payload["source_evidence_cards"], gallery)

    def test_text_provider_override_is_request_scoped(self):
        with (
            patch.object(core_config, "AI_TEXT_PROVIDER", "openai"),
            patch.object(core_config, "GEMINI_FOCUSED_MODEL", "gemini-focused"),
            patch.dict("os.environ", {"OPENAI_FOCUSED_MODEL": "openai-focused"}),
        ):
            self.assertEqual(core_config.active_text_provider(), "openai")
            self.assertEqual(core_config.model_for_depth("focused"), "openai-focused")
            token = core_config.set_request_text_provider("gemini")
            try:
                self.assertEqual(core_config.active_text_provider(), "gemini")
                self.assertEqual(core_config.model_for_depth("focused"), "gemini-focused")
            finally:
                core_config.reset_request_text_provider(token)
            self.assertEqual(core_config.active_text_provider(), "openai")


class VisualGalleryTests(unittest.TestCase):
    def test_pptx_svg_visual_rendering_respects_requested_slide_limit(self):
        slide_texts = [
            "Slide 1 table data comparison",
            "Slide 2 chart results",
            "Slide 3 agenda",
            "Slide 4 references",
        ]
        rendered = {
            1: "data:image/png;base64,AA==",
            2: "data:image/png;base64,AA==",
            3: "data:image/png;base64,AA==",
            4: "data:image/png;base64,AA==",
        }

        with (
            patch.object(backend_app_module, "ENABLE_PPTX_SLIDE_RENDER", False, create=True),
            patch.object(backend_app_module, "ENABLE_PPTX_SVG_FALLBACK_RENDER", True, create=True),
            patch.object(backend_app_module, "CONTROLLED_MAX_PPTX_SLIDES_PER_SOURCE", 6, create=True),
            patch.object(backend_app_module, "MAX_VISUAL_IMAGES_PER_SOURCE", 2, create=True),
            patch.object(backend_app_module, "Presentation", return_value=SimpleNamespace(slides=[]), create=True),
            patch.object(backend_app_module, "render_pptx_source_preview_svg_images", return_value=rendered),
            patch.object(backend_app_module, "selected_indices_by_score", return_value=[0, 1, 2, 3]),
        ):
            parts = backend_app_module.render_pptx_slide_screenshots(
                b"pptx bytes",
                "lecture.pptx",
                slide_texts,
                max_slides=2,
            )

        self.assertEqual(len(parts), 4)
        self.assertIn("PPT slide 1", parts[0]["text"])
        self.assertIn("PPT slide 2", parts[2]["text"])
        self.assertNotIn("PPT slide 3", json.dumps(parts))

    def test_pdf_visual_rendering_can_be_disabled_for_a_constrained_host(self):
        with (
            patch.object(backend_app_module, "ENABLE_PDF_VISUAL_EXTRACTION", False, create=True),
            patch("backend.app.extract_pdf", return_value="Lecture text"),
            patch("backend.app.render_pdf_visual_parts", side_effect=AssertionError("PDF rendering should be skipped")),
        ):
            _, source = file_to_source_unit("lecture.pdf", "application/pdf", b"pdf bytes")

        self.assertEqual(source["text_excerpt"], "Lecture text")
        self.assertEqual(source["visual_parts"], [])

    def test_uploaded_image_source_becomes_inline_visual_candidate_without_model_call(self):
        png_data = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
            "AAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
        )
        _, source = file_to_source_unit("results-chart.png", "image/png", png_data)
        source_units = [source]

        self.assertEqual(source["visual_parts"][0]["type"], "text")
        candidates = iter_visual_candidates(source_units)
        self.assertEqual(len(candidates), 1)
        self.assertFalse(candidates[0]["is_likely_decorative"])
        self.assertIn(candidates[0]["visual_kind"], {"data/table", "graph/chart", "diagram/model", "method/result figure"})

        with patch("backend.app.generate_chat", side_effect=AssertionError("uploaded image fallback should not call the model")):
            cards = rebuild_cached_visual_argument_cards(source_units, "english")
            summary = finalize_generated_summary(
                "## Notes\n\nThe uploaded results chart should be used as source evidence.",
                requested_language="english",
                generation_language="english",
                source_context="",
                source_units=source_units,
                attach_visuals=True,
            )
            gallery = build_visual_gallery(source_units)

        self.assertEqual(len(cards), 1)
        self.assertIn("[[VISUAL:0]]", summary)
        self.assertEqual(len(gallery), 1)
        assert_served_visual_asset(self, gallery[0]["url"], "image/png")

    def test_selected_visual_card_uses_marker_index(self):
        source_units = [{
            "visual_argument_cards": [{
                "index": 17,
                "url": "data:image/png;base64,AA==",
                "title": "Result comparison",
                "caption": "Table with data, results, comparison groups, and statistical evidence.",
                "what_shows": "The table compares results and shows a data pattern.",
                "argument_supported": "The table supports the analysis by showing the comparison directly.",
                "how_to_read": "Read rows, columns, values, and group differences.",
                "visual_kind": "unknown",
            }]
        }]

        gallery = build_visual_gallery(source_units)

        self.assertEqual(len(gallery), 1)
        self.assertEqual(gallery[0]["index"], 0)
        assert_served_visual_asset(self, gallery[0]["url"], "image/png")

    def test_cached_visual_rebuild_does_not_call_model(self):
        source_units = [{
            "display_name": "lecture.pdf",
            "title_candidate": "Lecture data",
            "visual_parts": [
                {
                    "type": "text",
                    "text": (
                        "IN-TEXT SOURCE FIGURE FROM lecture.pdf — PDF page 2. "
                        "Page text preview: table graph data results comparison. "
                        "Image-count=1; drawing-count=12; visual-score=28."
                    ),
                },
                {"type": "image_url", "image_url": {"url": "data:image/png;base64,AA=="}},
            ],
        }]

        with patch("backend.app.generate_chat", side_effect=AssertionError("model should not be called")):
            cards = rebuild_cached_visual_argument_cards(source_units, "english")
            summary = finalize_generated_summary(
                "## Notes\n\nThe table compares the source results.",
                requested_language="english",
                generation_language="english",
                source_context="",
                source_units=source_units,
                attach_visuals=True,
            )
            gallery = build_visual_gallery(source_units)

        self.assertEqual(len(cards), 1)
        self.assertIn("[[VISUAL:0]]", summary)
        self.assertNotIn("This source figure belongs in the notes", summary)
        self.assertNotIn("Figure focus:", summary)
        self.assertEqual(len(gallery), 1)
        assert_served_visual_asset(self, gallery[0]["url"], "image/png")

    def test_finalize_keeps_visual_cards_out_of_main_summary_and_is_idempotent(self):
        source_units = [{
            "visual_argument_cards": [{
                "index": 0,
                "url": "data:image/png;base64,AA==",
                "title": "Result table",
                "caption": "Table with data and results.",
                "what_shows": "The table compares results and shows a data pattern.",
                "argument_supported": "The table supports the analysis by showing the comparison directly.",
                "how_to_read": "Read rows, columns, values, and group differences.",
                "visual_kind": "data/table",
            }]
        }]
        polluted = (
            "# Generated Study Notes\n\n"
            "## Overview\n\n"
            "This lecture teaches how to interpret the result table as evidence.\n\n"
            "## Source Examples and Evidence\n\n"
            "### Result table\n\n"
            "This source figure belongs in the notes because it shows: Result table. "
            "It supports the surrounding explanation by making the method, pattern, mechanism, or contrast visible "
            "rather than asking the student to memorise an abstract claim.\n\n"
            "*Figure focus: Result table. It supports the surrounding explanation by making the method, pattern, mechanism, or contrast visible.*\n\n"
            "[[VISUAL:0]]\n\n"
            "## Core Argument\n\n"
            "The real study task is to explain what the evidence can and cannot prove."
        )

        once = finalize_generated_summary(
            polluted,
            requested_language="english",
            generation_language="english",
            source_context="",
            source_units=source_units,
            attach_visuals=True,
        )
        twice = finalize_generated_summary(
            once,
            requested_language="english",
            generation_language="english",
            source_context="",
            source_units=source_units,
            attach_visuals=True,
        )
        gallery = build_visual_gallery(source_units)

        self.assertEqual(once, twice)
        for bad_phrase in (
            "Source Examples and Evidence",
            "This source figure belongs in the notes",
            "method, pattern, mechanism, or contrast visible",
            "Figure focus:",
        ):
            self.assertNotIn(bad_phrase, twice)
        self.assertEqual(twice.count("[[VISUAL:0]]"), 1)
        self.assertIn("## Overview", twice)
        self.assertIn("## Core Argument", twice)
        self.assertIn("The real study task is to explain", twice)
        self.assertEqual(len(gallery), 1)

    def test_visual_gallery_serves_runtime_asset_url(self):
        source_units = [{
            "visual_argument_cards": [{
                "index": 0,
                "url": "data:image/png;base64,AA==",
                "title": "Source data table",
                "caption": "Table with data, results, comparison groups, and statistical evidence.",
                "what_shows": "The table compares results and shows a data pattern.",
                "argument_supported": "The table supports the analysis by showing the comparison directly.",
                "how_to_read": "Read rows, columns, values, and group differences.",
                "visual_kind": "data/table",
            }]
        }]

        gallery = build_visual_gallery(source_units)

        self.assertEqual(len(gallery), 1)
        self.assertFalse(gallery[0]["url"].startswith("data:"))
        assert_served_visual_asset(self, gallery[0]["url"], "image/png")

    def test_visual_asset_mirrors_new_data_url_to_durable_storage(self):
        png_data = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
            "AAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
        )
        data_url = "data:image/png;base64," + base64.b64encode(png_data).decode("ascii")

        with (
            patch.object(visual_assets, "durable_visual_storage_enabled", return_value=True),
            patch.object(visual_assets, "persist_visual_asset_to_durable_storage", return_value=True) as persist,
        ):
            url = visual_asset_url_for_browser(data_url)

        self.assertTrue(url.endswith(".png"))
        persist.assert_called_once()
        asset_name, content_type, stored_data = persist.call_args.args
        self.assertTrue(asset_name.endswith(".png"))
        self.assertEqual(content_type, "image/png")
        self.assertEqual(stored_data, png_data)

    def test_visual_asset_route_recovers_missing_runtime_file_from_durable_storage(self):
        asset_name = "durable-test-visual.png"
        asset_path = backend_app_module.RUNTIME_ASSETS_DIR / "visuals" / asset_name
        asset_path.unlink(missing_ok=True)
        png_data = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
            "AAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
        )

        try:
            with patch.object(
                backend_app_module,
                "fetch_visual_asset_from_durable_storage",
                return_value=(png_data, "image/png"),
            ) as fetch:
                response = TestClient(app).get(f"/assets/visuals/{asset_name}")
        finally:
            asset_path.unlink(missing_ok=True)

        self.assertEqual(response.status_code, 200)
        self.assertEqual(response.headers.get("content-type"), "image/png")
        self.assertEqual(response.content, png_data)
        fetch.assert_called_once_with(asset_name)

    def test_source_preview_image_url_serves_runtime_asset_when_requested(self):
        png_data = base64.b64decode(
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJ"
            "AAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
        )

        url = source_preview_image_url(png_data, "image/png", browser_asset=True)

        self.assertFalse(url.startswith("data:"))
        assert_served_visual_asset(self, url, "image/png")

    def test_visual_asset_url_preserves_existing_browser_url(self):
        url = "http://127.0.0.1:8001/assets/visuals/already.png"

        self.assertEqual(visual_asset_url_for_browser(url), url)

    def test_visual_asset_url_rejects_invalid_or_non_raster_data_urls(self):
        invalid_url = "data:image/png;base64,not-valid-base64"
        svg_url = "data:image/svg+xml;base64,PHN2Zy8+"

        self.assertEqual(visual_asset_url_for_browser(invalid_url), "")
        self.assertEqual(visual_asset_url_for_browser(svg_url), "")

    def test_visual_asset_url_rejects_local_filesystem_paths(self):
        self.assertEqual(visual_asset_url_for_browser("file:///Users/student/source/page.png"), "")
        self.assertEqual(visual_asset_url_for_browser("/Users/student/source/page.png"), "")
        self.assertEqual(visual_asset_url_for_browser("C:\\Users\\student\\source\\page.png"), "")

    def test_visual_gallery_rejects_local_filesystem_card_urls(self):
        source_units = [{
            "visual_argument_cards": [{
                "index": 0,
                "url": "/Users/student/source/page.png",
                "title": "Local-only table",
                "caption": "Table with data, results, comparison groups, and statistical evidence.",
                "what_shows": "The table compares results and shows a data pattern.",
                "argument_supported": "The table supports the analysis by showing the comparison directly.",
                "how_to_read": "Read rows, columns, values, and group differences.",
                "visual_kind": "data/table",
            }]
        }]

        gallery = build_visual_gallery(source_units)

        self.assertEqual(gallery, [])

    def test_cached_id_only_visual_metadata_preserves_matching_marker(self):
        summary = "## Notes\n\nUse the evidence table here.\n\n[[VISUAL:7]]\n\nNext point."
        items = [{
            "id": 7,
            "url": "http://127.0.0.1:8001/assets/visuals/id-only.png",
            "title": "ID-only source figure",
        }]

        pruned = prune_unavailable_visual_markers(summary, items)

        self.assertIn("[[VISUAL:7]]", pruned)

    def test_visual_markers_are_reindexed_after_unrenderable_cards_are_removed(self):
        source_units = [{
            "visual_argument_cards": [
                {
                    "index": 0,
                    "url": "",
                    "title": "Missing visual",
                    "caption": "Table with data, results, comparison groups, and statistical evidence.",
                    "what_shows": "The unavailable table compares results.",
                    "argument_supported": "This card should not get a marker because it has no image URL.",
                    "how_to_read": "Read rows, columns, values, and group differences.",
                    "visual_kind": "data/table",
                },
                {
                    "index": 1,
                    "url": "data:image/png;base64,AA==",
                    "title": "Renderable result table",
                    "caption": "Table with data, results, comparison groups, and statistical evidence.",
                    "what_shows": "The table compares results and shows a data pattern.",
                    "argument_supported": "The table supports the analysis by showing the comparison directly.",
                    "how_to_read": "Read rows, columns, values, and group differences.",
                    "visual_kind": "data/table",
                },
            ]
        }]

        summary = finalize_generated_summary(
            "## Notes\n\nThe table compares source results and supports the key claim.",
            requested_language="english",
            generation_language="english",
            source_context="",
            source_units=source_units,
            attach_visuals=True,
        )
        gallery = build_visual_gallery(source_units)

        self.assertIn("[[VISUAL:0]]", summary)
        self.assertNotIn("[[VISUAL:1]]", summary)
        self.assertEqual(len(gallery), 1)
        self.assertEqual(gallery[0]["index"], 0)
        self.assertEqual(gallery[0]["title"], "Renderable result table")
        assert_served_visual_asset(self, gallery[0]["url"], "image/png")

    def test_pptx_svg_fallback_is_rasterized_for_runtime_assets(self):
        if Presentation is None:
            self.skipTest("python-pptx is not installed")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(914400, 914400, 5486400, 914400)
        textbox.text = "Results table and comparison graph"

        rendered = render_pptx_source_preview_svg_images(prs, 1)

        self.assertIn(1, rendered)
        self.assertTrue(rendered[1].startswith("data:image/png;base64,"))
        assert_served_visual_asset(self, visual_asset_url_for_browser(rendered[1]), "image/png")

    def test_pptx_svg_browser_preview_uses_runtime_asset_url(self):
        if Presentation is None:
            self.skipTest("python-pptx is not installed")
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        textbox = slide.shapes.add_textbox(914400, 914400, 5486400, 914400)
        textbox.text = "Results table and comparison graph"

        rendered = render_pptx_source_preview_svg_images(prs, 1, browser_assets=True)

        self.assertIn(1, rendered)
        self.assertFalse(rendered[1].startswith("data:"))
        assert_served_visual_asset(self, rendered[1], "image/png")

    def test_youtube_link_preserves_frame_visual_parts_for_inline_candidates(self):
        frame_parts = [
            {
                "type": "text",
                "text": (
                    "IN-TEXT SOURCE FIGURE FROM lecture video — video frame 1 sampled at approximately 00:10. "
                    "The frame shows a graph with data, results, comparison groups, and visual-score=32."
                ),
            },
            {"type": "image_url", "image_url": {"url": "data:image/jpeg;base64,AA=="}},
        ]
        meta = {
            "url": "https://www.youtube.com/watch?v=dQw4w9WgXcQ",
            "source_identity": "youtube:dQw4w9WgXcQ",
            "detected_title": "Lecture video",
            "content_hash": "hash",
            "visual_parts": frame_parts,
        }

        with patch("backend.app.analyse_youtube_url", return_value=("Transcript text", frame_parts, meta)):
            parts, source = link_to_source_unit("https://www.youtube.com/watch?v=dQw4w9WgXcQ")

        self.assertEqual(source["visual_parts"], frame_parts)
        self.assertEqual(parts[-1]["image_url"]["url"], "data:image/jpeg;base64,AA==")
        candidates = iter_visual_candidates([source])
        self.assertEqual(len(candidates), 1)
        self.assertEqual(candidates[0]["url"], "data:image/jpeg;base64,AA==")


class SourceIdentityTests(unittest.TestCase):
    def test_uploaded_file_identity_keeps_raw_bytes_when_extracted_text_matches(self):
        with patch("backend.app.extract_pdf", return_value=""), patch("backend.app.render_pdf_visual_parts", return_value=[]):
            _, first = file_to_source_unit("scan-a.pdf", "application/pdf", b"pdf bytes a")
            _, second = file_to_source_unit("scan-b.pdf", "application/pdf", b"pdf bytes b")

        self.assertNotEqual(first["source_identity"], second["source_identity"])
        self.assertNotEqual(first["content_hash"], second["content_hash"])
        self.assertNotEqual(
            build_analysis_fingerprint("auto", [first], "detailed", "professor_mode"),
            build_analysis_fingerprint("auto", [second], "detailed", "professor_mode"),
        )


class AnalysisCacheTests(unittest.TestCase):
    def test_tight_analysis_budget_skips_optional_visual_model_filter(self):
        source_units = [{
            "display_name": "budget.pdf",
            "title_candidate": "Budget Control Lecture",
            "source_identity": "file:budget",
            "text_excerpt": "This lecture compares evidence tables, diagrams, and revision uses for budget control.",
            "visual_parts": [
                {
                    "type": "text",
                    "text": (
                        "IN-TEXT SOURCE FIGURE FROM budget.pdf — PDF page 2. "
                        "Actual source screenshot selected for its teaching figure/table/data value. "
                        "Image-count=1; drawing-count=22; visual-score=82. "
                        "Page text preview: Table comparing concepts, evidence, and study use."
                    ),
                },
                {
                    "type": "image_url",
                    "image_url": {
                        "url": (
                            "data:image/png;base64,"
                            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/p9sAAAAASUVORK5CYII="
                        ),
                    },
                },
            ],
        }]
        started_at = time.monotonic() - (backend_app_module.ANALYSIS_MAX_SECONDS - 100)
        skipped = []
        generated_calls = []

        def fake_generate_chat(messages, **kwargs):
            call_text = json.dumps(messages, default=str)
            generated_calls.append(call_text)
            self.assertNotIn("selecting separate source-figure support cards", call_text)
            return "# Budget-Controlled Study Notes\n\n" + (
                "This source-grounded note explains the concepts, evidence, diagram, comparison table, "
                "and revision use without spending an optional visual-filter model call. "
                * 80
            )

        with (
            patch("backend.app.generate_chat", side_effect=fake_generate_chat),
            patch("backend.app.advanced_notes_quality_flags", return_value=[]),
            patch("backend.app.markdown_table_count", return_value=3),
        ):
            result = generate_reference_style_multisource_notes(
                source_units,
                preferred_language="english",
                depth_plan={"depth": "detailed", "config": backend_app_module.DEPTH_CONFIG["detailed"]},
                prompt_mode="professor_mode",
                note_length_mode="standard_notes",
                analysis_started_at=started_at,
                skipped_optional_stages=skipped,
            )

        self.assertIn("Budget-Controlled Study Notes", result)
        self.assertEqual(len(generated_calls), 1)
        self.assertIn("visual_card_filter", skipped)

    def test_cache_round_trip_and_ttl_cleanup(self):
        with tempfile.TemporaryDirectory() as tmp:
            cache_path = Path(tmp) / "cache.json"
            store = AnalysisCacheStore(cache_path=cache_path, ttl_seconds=60)

            store.set("fingerprint", {"summary": "ok", "visual_gallery": []})

            self.assertEqual(store.get("fingerprint")["summary"], "ok")
            self.assertTrue(cache_path.exists())

            cache_path.write_text(
                json.dumps({
                    "expired": {
                        "created_at": time.time() - 120,
                        "result": {"summary": "old"},
                    }
                }),
                encoding="utf-8",
            )
            expired_store = AnalysisCacheStore(cache_path=cache_path, ttl_seconds=1)

            self.assertIsNone(expired_store.get("expired"))

    def test_cache_save_uses_complete_file_without_temp_artifacts(self):
        with tempfile.TemporaryDirectory() as tmp:
            cache_path = Path(tmp) / "cache.json"
            store = AnalysisCacheStore(cache_path=cache_path, ttl_seconds=60)

            store.set("fingerprint", {"summary": "ok", "visual_gallery": [{"index": 0}]})

            self.assertEqual(store.get("fingerprint")["summary"], "ok")
            self.assertFalse(list(Path(tmp).glob(".cache.json.*.tmp")))

    def test_analyze_cache_hit_does_not_require_text_ai_client(self):
        cached_result = {
            "title": "Cached Study Notes",
            "summary": "# Cached Study Notes\n\nCached body from an earlier successful analysis.",
            "display_summary": "# Cached Study Notes\n\nCached body from an earlier successful analysis.",
            "sections": {"Cached Study Notes": "Cached body from an earlier successful analysis."},
            "connections": [],
            "mind_map": {"center": "Cached Study Notes", "branches": []},
            "source_identity": "text:cached",
            "primary_source_identity": "text:cached",
            "visual_gallery": [],
            "visuals": [],
        }

        with (
            patch("backend.app.require_text_ai", side_effect=RuntimeError("OPENAI_API_KEY is not configured.")) as require_text_ai,
            patch("backend.app.cache_get", return_value=cached_result) as cache_get,
            patch("backend.app.generate_chat", side_effect=AssertionError("cache hit should not call the model")),
            patch("backend.app.persist_generated_analysis_result", return_value={}),
        ):
            result = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="Cached source text",
                preferred_language="english",
                detail_level="standard",
                prompt_mode="professor_mode",
                note_length="standard_notes",
                ai_provider="openai",
                client_fingerprint="cache-client",
                request=None,
            ))

        self.assertIsInstance(result, dict)
        self.assertTrue(result["cached"])
        self.assertEqual(result["title"], "Cached Study Notes")
        self.assertIn("Cached body", result["summary"])
        require_text_ai.assert_not_called()
        cache_get.assert_called_once()

    def test_analyze_records_safe_stage_diagnostics_for_cache_hit(self):
        cached_result = {
            "title": "Cached Study Notes",
            "summary": "# Cached Study Notes\n\nCached body.",
            "display_summary": "# Cached Study Notes\n\nCached body.",
            "sections": {"Cached Study Notes": "Cached body."},
            "connections": [],
            "mind_map": {"center": "Cached Study Notes", "branches": []},
            "source_identity": "text:cached",
            "primary_source_identity": "text:cached",
            "visual_gallery": [],
            "visuals": [],
        }
        emitted_events = []

        def capture_event(message, *args, **kwargs):
            emitted_events.append(message % args if args else message)

        with (
            patch("backend.app.cache_get", return_value=cached_result),
            patch("backend.app.persist_generated_analysis_result", return_value={}),
            patch("backend.app.logger.info", side_effect=capture_event),
        ):
            result = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="A short source about opportunity cost.",
                preferred_language="english",
                detail_level="standard",
                prompt_mode="professor_mode",
                note_length="standard_notes",
                ai_provider="openai",
                client_fingerprint="diagnostic-client",
                request=None,
            ))

        self.assertTrue(result["cached"])
        self.assertTrue(any("analysis_event=received" in event for event in emitted_events))
        self.assertTrue(any("analysis_event=sources_ready" in event for event in emitted_events))
        self.assertTrue(any("analysis_event=cache_hit" in event for event in emitted_events))

    def test_analyze_stops_before_model_generation_when_client_disconnected(self):
        class DisconnectedRequest:
            async def is_disconnected(self):
                return True

        with (
            patch("backend.app.cache_get", return_value=None),
            patch("backend.app.require_text_ai") as require_text_ai,
            patch(
                "backend.app.generate_reference_style_multisource_notes",
                side_effect=AssertionError("model should not be called after disconnect"),
            ) as generate_notes,
            patch("backend.app.generate_chat", side_effect=AssertionError("model should not be called after disconnect")),
            patch("backend.app.persist_generated_analysis_result", return_value={}),
        ):
            response = asyncio.run(analyze_materials(
                files=[],
                links="[]",
                free_text="A short source about opportunity cost and scarce resources.",
                preferred_language="english",
                detail_level="focused",
                prompt_mode="professor_mode",
                note_length="standard_notes",
                ai_provider="openai",
                client_fingerprint="disconnect-client",
                request=DisconnectedRequest(),
            ))

        self.assertEqual(response.status_code, 499)
        self.assertIn("application/json", response.media_type)
        self.assertIn("Client disconnected", response.body.decode("utf-8"))
        require_text_ai.assert_not_called()
        generate_notes.assert_not_called()


if __name__ == "__main__":
    unittest.main()
